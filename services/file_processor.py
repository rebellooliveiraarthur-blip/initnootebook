from pathlib import Path
from io import BytesIO
import PyPDF2
from docx import Document
import pptx

class FileProcessor:
    def process(self, uploaded_file, file_bytes):
        """Identify file type and extract text"""
        file_extension = Path(uploaded_file.name).suffix.lower()
        
        if file_extension == '.txt':
            return self._process_txt(file_bytes)
        elif file_extension == '.pdf':
            return self._process_pdf(file_bytes)
        elif file_extension == '.docx':
            return self._process_docx(file_bytes)
        elif file_extension == '.pptx':
            return self._process_pptx(file_bytes)
        else:
            raise ValueError(f"Unsupported file type: {file_extension}")
    
    def _process_txt(self, file_bytes):
        return file_bytes.decode('utf-8')
    
    def _process_pdf(self, file_bytes):
        pdf_reader = PyPDF2.PdfReader(BytesIO(file_bytes))
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text()
        return text
    
    def _process_docx(self, file_bytes):
        doc = Document(BytesIO(file_bytes))
        text = "\n".join([p.text for p in doc.paragraphs])
        return text
    
    def _process_pptx(self, file_bytes):
        prs = pptx.Presentation(BytesIO(file_bytes))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text